import os
import tempfile
import subprocess
from pathlib import Path
import re
import time

def clean_text(text):
    """Nettoie le texte extrait."""
    # Supprimer les caractères spéciaux et les retours à la ligne multiples
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[^\w\s.,;:!?()-]', '', text)
    return text

def extract_text_from_pptx(file_path):
    """Extrait le texte d'un fichier PowerPoint (.pptx ou .ppt)."""
    try:
        # Vérifier que le fichier existe
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"Le fichier {file_path} n'existe pas.")
        
        # Créer un fichier temporaire pour la sortie
        with tempfile.NamedTemporaryFile(suffix='.txt', delete=False) as tmp_file:
            output_file = tmp_file.name
        
        # Utiliser un outil externe comme LibreOffice pour convertir PPT en texte
        # Cette approche nécessite que LibreOffice soit installé
        try:
            cmd = ['soffice', '--headless', '--convert-to', 'txt', file_path, '--outdir', os.path.dirname(output_file)]
            process = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            stdout, stderr = process.communicate()
            
            # Attendre la fin de la conversion
            if process.returncode != 0:
                raise Exception(f"Erreur lors de la conversion du fichier: {stderr.decode('utf-8')}")
            
            # Trouver le fichier de sortie
            base_name = os.path.basename(file_path)
            name_without_ext = os.path.splitext(base_name)[0]
            converted_file = os.path.join(os.path.dirname(output_file), f"{name_without_ext}.txt")
            
            # Attendre que le fichier soit créé
            timeout = 10  # secondes
            start_time = time.time()
            while not os.path.exists(converted_file) and time.time() - start_time < timeout:
                time.sleep(0.5)
            
            if not os.path.exists(converted_file):
                raise FileNotFoundError(f"Le fichier converti n'a pas été trouvé après {timeout} secondes.")
            
            # Lire le contenu du fichier
            with open(converted_file, 'r', encoding='utf-8', errors='ignore') as f:
                text = f.read()
            
            # Nettoyer le texte
            text = clean_text(text)
            
            # Supprimer les fichiers temporaires
            if os.path.exists(converted_file):
                os.remove(converted_file)
            if os.path.exists(output_file):
                os.remove(output_file)
            
            return text
            
        except (FileNotFoundError, Exception):
            # Utiliser une approche alternative avec python-pptx
            import pptx
            presentation = pptx.Presentation(file_path)
            text_list = []
            
            # Extraire le texte de chaque slide
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_list.append(shape.text)
            
            text = "\n\n".join(text_list)
            return clean_text(text)
    
    except Exception as e:
        print(f"Erreur lors de l'extraction du texte du fichier PPT: {str(e)}")
        return ""

class PPTXTextLoader:
    """Chargeur personnalisé pour les fichiers PowerPoint."""
    
    def __init__(self, file_path):
        self.file_path = file_path
        
    def load(self):
        """Charge le contenu du fichier PowerPoint comme un document LangChain."""
        from langchain.schema import Document
        
        text = extract_text_from_pptx(self.file_path)
        metadata = {"source": self.file_path}
        
        return [Document(page_content=text, metadata=metadata)]
